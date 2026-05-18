#!/usr/bin/env python3
"""
Generate Architecture PowerPoint Presentation
for ajaib-testing-code project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Ajaib Testing Code"
    subtitle.text = "Architecture Overview\nHexagonal Architecture Pattern\n\nTransfer Service Implementation"

    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_overview_slide(prs):
    """Create architecture overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Overview"

    content = slide.placeholders[1].text_frame
    content.text = "Hexagonal Architecture (Ports and Adapters Pattern)"

    points = [
        "Clean separation of concerns",
        "Business logic independent of external frameworks",
        "Testable and maintainable codebase",
        "Technology-agnostic core domain",
        "Easy to swap implementations (DB, Cache, HTTP framework)"
    ]

    for point in points:
        p = content.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)

def create_hexagonal_architecture_slide(prs):
    """Create hexagonal architecture explanation slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Hexagonal Architecture Pattern"

    content = slide.placeholders[1].text_frame
    content.text = "Core Concepts:"

    concepts = [
        ("Ports", "Interfaces that define contracts"),
        ("Adapters", "Implementations of port interfaces"),
        ("Primary/Driving Adapters", "Trigger actions (HTTP handlers, CLI)"),
        ("Secondary/Driven Adapters", "Provide data/services (DB, Cache, APIs)"),
        ("Domain/Core", "Business logic, independent of infrastructure")
    ]

    for concept, description in concepts:
        p = content.add_paragraph()
        p.text = f"{concept}: {description}"
        p.level = 1
        p.font.size = Pt(16)

def create_layers_slide(prs):
    """Create layers breakdown slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Layers"

    content = slide.placeholders[1].text_frame
    content.text = "Layer Structure:"

    layers = [
        ("1. Entity Layer", "Domain models (Transfer, CreateTransferRequest)"),
        ("2. Ports Layer", "Interface definitions"),
        ("   - ports/app", "Application service interfaces"),
        ("   - ports/core", "Core business logic interfaces"),
        ("   - ports/secondary", "Secondary adapter interfaces (DB, Cache)"),
        ("3. Adapters Layer", "Interface implementations"),
        ("   - adapters/core", "Business logic implementation"),
        ("   - adapters/app", "Application service implementation"),
        ("   - adapters/framework/primary", "REST handlers (Gin)"),
        ("   - adapters/framework/secondary", "DB & Cache repositories"),
        ("4. Router", "HTTP routing configuration"),
        ("5. CMD", "Application entry point (main.go)")
    ]

    for layer, description in layers:
        p = content.add_paragraph()
        p.text = f"{layer}: {description}"
        p.level = 0 if layer[0].isdigit() else 1
        p.font.size = Pt(14)

def create_data_flow_slide(prs):
    """Create data flow slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Request Flow"

    content = slide.placeholders[1].text_frame
    content.text = "HTTP Request → Response Flow:"
    content.paragraphs[0].font.bold = True

    flow_steps = [
        "1. HTTP Request arrives at Gin Router",
        "2. Router dispatches to Handler (Primary Adapter)",
        "3. Handler validates request and calls App Service (Port)",
        "4. App Service orchestrates business flow",
        "5. App Service calls Core Service (Port)",
        "6. Core Service executes business logic",
        "7. Core Service calls Secondary Adapters (DB/Cache)",
        "8. Secondary Adapters interact with infrastructure",
        "9. Response flows back through layers",
        "10. Handler returns HTTP response"
    ]

    for step in flow_steps:
        p = content.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(16)

def create_tech_stack_slide(prs):
    """Create technology stack slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"

    content = slide.placeholders[1].text_frame
    content.text = "Core Technologies:"

    tech_items = [
        ("Language", "Go 1.25.5"),
        ("HTTP Framework", "Gin (github.com/gin-gonic/gin)"),
        ("Testing", "Go standard testing + go.uber.org/mock"),
        ("Logging", "slog (structured logging)"),
        ("Architecture", "Hexagonal (Ports & Adapters)"),
        ("Dependency Injection", "Constructor-based DI")
    ]

    for tech, detail in tech_items:
        p = content.add_paragraph()
        p.text = f"{tech}: {detail}"
        p.level = 1
        p.font.size = Pt(18)

def create_directory_structure_slide(prs):
    """Create directory structure slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Directory Structure"

    content = slide.placeholders[1].text_frame
    content.text = "Project Organization:"

    structure = [
        "cmd/gateway/ - Application entry point",
        "config/ - Configuration management",
        "internal/ - Private application code",
        "  adapters/ - Port implementations",
        "    app/ - Application services",
        "    core/ - Business logic",
        "      entity/ - Domain models",
        "    framework/ - Infrastructure",
        "      primary/ - Driving adapters (REST)",
        "      secondary/ - Driven adapters (DB, Cache)",
        "  ports/ - Interface definitions",
        "    app/ - App service interfaces",
        "    core/ - Core service interfaces",
        "    secondary/ - Infrastructure interfaces",
        "router/ - HTTP routing setup",
        "test/ - Test suites",
        "  e2e/ - End-to-end tests",
        "  integration/ - Integration tests",
        "  mocks/ - Generated mocks"
    ]

    for item in structure:
        p = content.add_paragraph()
        p.text = item
        p.level = 0 if not item.startswith("  ") else (1 if item.startswith("  ") and not item.startswith("    ") else 2)
        p.font.size = Pt(13)

def create_testing_strategy_slide(prs):
    """Create testing strategy slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Testing Strategy"

    content = slide.placeholders[1].text_frame
    content.text = "Multi-Layer Testing Approach:"

    testing_layers = [
        ("Unit Tests", "Test individual components in isolation"),
        ("  - Handler tests", "Test HTTP handlers with mocked services"),
        ("  - Service tests", "Test business logic with mocked repositories"),
        ("  - Repository tests", "Test data access layer"),
        ("Integration Tests", "Test component interactions"),
        ("  - Service + Repository", "Test business logic with real dependencies"),
        ("E2E Tests", "Test complete request flow"),
        ("  - Full stack testing", "HTTP request → response validation"),
        ("Mock Generation", "Using go.uber.org/mock for test doubles")
    ]

    for layer, description in testing_layers:
        p = content.add_paragraph()
        p.text = f"{layer}: {description}"
        p.level = 0 if not layer.startswith("  ") else 1
        p.font.size = Pt(15)

def create_transfer_service_slide(prs):
    """Create transfer service details slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Transfer Service - Key Features"

    content = slide.placeholders[1].text_frame
    content.text = "Core Functionality:"

    features = [
        "Create Transfer - Initiate money transfer between accounts",
        "Get Transfer by ID - Retrieve transfer details",
        "List Transfers - Get all transfers",
        "Update Transfer Status - Modify transfer state",
        "",
        "Key Components:",
        "  • Idempotency handling via cache",
        "  • Balance tracking (from/to accounts)",
        "  • Status management",
        "  • Structured logging with slog",
        "  • Request validation",
        "  • Error handling and HTTP status codes"
    ]

    for feature in features:
        if feature == "":
            content.add_paragraph()
            continue
        p = content.add_paragraph()
        p.text = feature
        p.level = 0 if not feature.startswith("  ") else 1
        p.font.size = Pt(16)

def create_dependency_injection_slide(prs):
    """Create dependency injection slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Dependency Injection Flow"

    content = slide.placeholders[1].text_frame
    content.text = "Constructor-Based DI (from main.go):"

    di_flow = [
        "1. Initialize DB Repository",
        "   dbRepo := transferDB.New(config)",
        "",
        "2. Initialize Cache Repository",
        "   cacheRepo := idempotencyCache.New(config)",
        "",
        "3. Initialize Core Service (inject repositories)",
        "   coreTransfer := transferCore.New(dbRepo, cacheRepo)",
        "",
        "4. Initialize App Service (inject core)",
        "   appTransfer := transferApp.New(coreTransfer)",
        "",
        "5. Initialize Handler (inject app service)",
        "   handler := transferHandler.New(appTransfer)",
        "",
        "6. Initialize Router (inject handler)",
        "   router := router.NewRouter(handler)"
    ]

    for step in di_flow:
        if step == "":
            content.add_paragraph()
            continue
        p = content.add_paragraph()
        p.text = step
        p.level = 0 if step[0].isdigit() else 1
        p.font.size = Pt(14)

def create_benefits_slide(prs):
    """Create architecture benefits slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Benefits"

    content = slide.placeholders[1].text_frame
    content.text = "Why Hexagonal Architecture?"

    benefits = [
        ("Testability", "Easy to mock dependencies, test in isolation"),
        ("Maintainability", "Clear separation makes changes localized"),
        ("Flexibility", "Swap implementations without changing core logic"),
        ("Independence", "Business logic doesn't depend on frameworks"),
        ("Scalability", "Easy to add new adapters (gRPC, GraphQL, etc.)"),
        ("Team Collaboration", "Clear boundaries enable parallel development"),
        ("Technology Agnostic", "Core domain independent of tech choices")
    ]

    for benefit, description in benefits:
        p = content.add_paragraph()
        p.text = f"{benefit}: {description}"
        p.level = 1
        p.font.size = Pt(16)

def main():
    """Main function to generate PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating slides...")
    create_title_slide(prs)
    create_overview_slide(prs)
    create_hexagonal_architecture_slide(prs)
    create_layers_slide(prs)
    create_data_flow_slide(prs)
    create_tech_stack_slide(prs)
    create_directory_structure_slide(prs)
    create_testing_strategy_slide(prs)
    create_transfer_service_slide(prs)
    create_dependency_injection_slide(prs)
    create_benefits_slide(prs)

    output_file = "/home/huzairuje/go/src/ajaib-testing-code/Architecture_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
